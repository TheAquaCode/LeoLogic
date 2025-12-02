"""
Multi-Modal File Extractors
Supports: Documents, Images, Audio, Video, Archives, Adobe files
"""

from pathlib import Path
from typing import Dict, List, Any
import json
import mimetypes
import PyPDF2
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
from PIL import Image
import pytesseract
from faster_whisper import WhisperModel
import subprocess

try:
    from psd_tools import PSDImage

    HAS_PSD = True
except ImportError:
    HAS_PSD = False
import zipfile
from .state import state
from utils.logger import setup_logger

logger = setup_logger(__name__)


class MultiModalExtractor:
    """Extract content from any file type"""

    @staticmethod
    def extract(file_path: str) -> Dict[str, Any]:
        """
        Extract all content from a file.
        Returns: {
            "text": str,
            "images": List[str],  # Base64 or paths
            "audio_transcript": str,
            "metadata": dict,
            "file_type": str
        }
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        result = {
            "text": "",
            "images": [],
            "audio_transcript": "",
            "metadata": {
                "filename": path.name,
                "size": path.stat().st_size,
                "extension": ext,
            },
            "file_type": MultiModalExtractor._detect_type(ext),
        }

        try:
            if ext in [
                ".txt",
                ".md",
                ".json",
                ".xml",
                ".html",
                ".css",
                ".js",
                ".py",
                ".java",
                ".cpp",
                ".csv",
            ]:
                result["text"] = MultiModalExtractor._extract_text(file_path)

            elif ext == ".pdf":
                result.update(MultiModalExtractor._extract_pdf(file_path))

            elif ext in [".docx", ".doc"]:
                result.update(MultiModalExtractor._extract_docx(file_path))

            elif ext in [".xlsx", ".xls"]:
                result["text"] = MultiModalExtractor._extract_excel(file_path)

            elif ext in [".pptx", ".ppt"]:
                result.update(MultiModalExtractor._extract_pptx(file_path))

            elif ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"]:
                result.update(MultiModalExtractor._extract_image(file_path))

            elif ext == ".psd" and HAS_PSD:
                result.update(MultiModalExtractor._extract_psd(file_path))

            elif ext in [".mp3", ".wav", ".m4a", ".flac", ".ogg", ".wma"]:
                result["audio_transcript"] = MultiModalExtractor._extract_audio(
                    file_path
                )

            elif ext in [".mp4", ".avi", ".mov", ".mkv", ".webm"]:
                result["audio_transcript"] = MultiModalExtractor._extract_video_audio(
                    file_path
                )

            elif ext == ".svg":
                result.update(MultiModalExtractor._extract_svg(file_path))

            elif ext == ".zip":
                result["text"] = MultiModalExtractor._extract_zip(file_path)

            else:
                result["text"] = f"Filename: {path.stem}"

        except Exception as e:
            logger.warning(f"Error extracting {file_path}: {e}")
            result["text"] = path.stem
            result["metadata"]["error"] = str(e)

        return result

    @staticmethod
    def _detect_type(ext: str) -> str:
        text_types = [
            ".txt",
            ".md",
            ".json",
            ".xml",
            ".html",
            ".css",
            ".js",
            ".py",
            ".java",
            ".csv",
        ]
        doc_types = [".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt"]
        image_types = [
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".tiff",
            ".webp",
            ".psd",
            ".svg",
        ]
        audio_types = [".mp3", ".wav", ".m4a", ".flac", ".ogg", ".wma"]
        video_types = [".mp4", ".avi", ".mov", ".mkv", ".webm"]

        if ext in text_types or ext in doc_types:
            return "document"
        elif ext in image_types:
            return "image"
        elif ext in audio_types:
            return "audio"
        elif ext in video_types:
            return "video"
        else:
            return "other"

    @staticmethod
    def _extract_text(file_path: str) -> str:
        """Extract plain text"""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    def _extract_pdf(file_path: str) -> Dict:
        """Extract text and images from PDF"""
        text = ""
        images = []

        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
                text += "\n"

        return {"text": text, "images": images}

    @staticmethod
    def _extract_docx(file_path: str) -> Dict:
        """Extract text and images from Word"""
        doc = DocxDocument(file_path)
        text = "\n".join([p.text for p in doc.paragraphs])

        images = []

        return {"text": text, "images": images}

    @staticmethod
    def _extract_excel(file_path: str) -> str:
        """Extract data from Excel"""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"\n=== Sheet: {sheet_name} ===\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(
                    [str(cell) if cell is not None else "" for cell in row]
                )
                text += row_text + "\n"

        return text

    @staticmethod
    def _extract_pptx(file_path: str) -> Dict:
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        text = ""
        images = []

        for i, slide in enumerate(prs.slides):
            text += f"\n=== Slide {i+1} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return {"text": text, "images": images}

    @staticmethod
    def _extract_image(file_path: str) -> Dict:
        """Extract text from image using OCR"""
        img = Image.open(file_path)
        try:
            text = pytesseract.image_to_string(img)
        except Exception as e:
            logger.warning(f"OCR failed for {file_path}: {e}")
            text = ""
        images = [file_path]

        return {"text": text, "images": images}

    @staticmethod
    def _extract_psd(file_path: str) -> Dict:
        """Extract preview from Photoshop"""
        psd = PSDImage.open(file_path)
        text = "Layers: " + ", ".join([layer.name for layer in psd])
        images = []

        return {"text": text, "images": images}

    @staticmethod
    def _extract_audio(file_path: str) -> str:
        """Transcribe audio using Whisper"""
        if not state.whisper_model:
            return ""

        try:
            segments, info = state.whisper_model.transcribe(file_path, beam_size=5)
            transcript = " ".join([segment.text for segment in segments])
            return transcript
        except Exception as e:
            logger.error(f"Whisper error: {e}")
            return ""

    @staticmethod
    def _extract_video_audio(file_path: str) -> str:
        """Extract and transcribe audio from video"""
        import tempfile

        temp_audio = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
        temp_audio_path = temp_audio.name
        temp_audio.close()

        try:
            subprocess.run(
                [
                    "ffmpeg",
                    "-i",
                    file_path,
                    "-vn",
                    "-acodec",
                    "pcm_s16le",
                    "-ar",
                    "16000",
                    "-ac",
                    "1",
                    temp_audio_path,
                ],
                capture_output=True,
                check=True,
            )
            transcript = MultiModalExtractor._extract_audio(temp_audio_path)
            Path(temp_audio_path).unlink()

            return transcript
        except Exception as e:
            logger.error(f"Video audio extraction error: {e}")
            return ""

    @staticmethod
    def _extract_svg(file_path: str) -> Dict:
        """Extract text from SVG file"""
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                # SVG is XML, so we can extract text content
                return {
                    "text": f"SVG file: {Path(file_path).stem}",
                    "images": [file_path],
                }
        except Exception as e:
            logger.warning(f"Error extracting SVG {file_path}: {e}")
            return {"text": Path(file_path).stem, "images": []}

    @staticmethod
    def _extract_zip(file_path: str) -> str:
        """List contents of ZIP"""
        with zipfile.ZipFile(file_path, "r") as zip_ref:
            file_list = zip_ref.namelist()
            return "Archive contents:\n" + "\n".join(file_list)
