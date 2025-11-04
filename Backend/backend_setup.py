"""
Leologic Backend Setup Script - FIXED VERSION
==============================================
Complete script with all missing file definitions.

Usage:
    python setup_backend.py
"""

import os
from pathlib import Path


def create_file(path: Path, content: str):
    """Create a file with the given content"""
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"✅ Created: {path}")


def setup_backend():
    """Generate the entire backend structure"""

    backend_dir = Path(__file__).parent

    print("🚀 Setting up Leologic Backend with Local AI + RAG...")
    print(f"📁 Target directory: {backend_dir}")
    print()

    # ============================================
    # MAIN.PY
    # ============================================

    main_py = '''"""
Leologic Backend - Main Entry Point
===================================
Boots up file organizer, chatbot, and RAG services.
"""

import os
import sys
from threading import Thread
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from file_organizer.app import create_file_organizer_app
from chatbot.app import create_chatbot_app
from config.settings import FILE_ORGANIZER_PORT, CHATBOT_PORT


def run_file_organizer():
    """Run file organizer service on port 5001"""
    print(f"🗂️  Starting File Organizer on port {FILE_ORGANIZER_PORT}...")
    app = create_file_organizer_app()
    app.run(host="0.0.0.0", port=FILE_ORGANIZER_PORT, debug=False, use_reloader=False)


def run_chatbot():
    """Run chatbot service on port 5000"""
    print(f"🤖 Starting Chatbot on port {CHATBOT_PORT}...")
    app = create_chatbot_app()
    app.run(host="0.0.0.0", port=CHATBOT_PORT, debug=False, use_reloader=False)


if __name__ == "__main__":
    print("=" * 60)
    print("🚀 LEOLOGIC - AI File Sorter Backend")
    print("=" * 60)
    print()
    
    # Start file organizer in separate thread
    organizer_thread = Thread(target=run_file_organizer, daemon=True)
    organizer_thread.start()
    
    # Start chatbot in main thread
    try:
        run_chatbot()
    except KeyboardInterrupt:
        print("\\n\\n👋 Shutting down Leologic Backend...")
        sys.exit(0)
'''

    # ============================================
    # CONFIG FILES
    # ============================================

    config_init = """# Config package"""

    config_settings = '''"""
Configuration Settings
"""

import os
from pathlib import Path
from dotenv import load_dotenv

load_dotenv()

# Base directories
BASE_DIR = Path(__file__).parent.parent
DATA_DIR = BASE_DIR / "data"
DATA_DIR.mkdir(exist_ok=True)

# RAG storage
RAG_DIR = DATA_DIR / "rag_store"
RAG_DIR.mkdir(exist_ok=True)

# Ports
FILE_ORGANIZER_PORT = 5001
CHATBOT_PORT = 5000

# File Organizer Settings
CONFIG_FILE = DATA_DIR / "file_organizer_config.json"
CONFIDENCE_THRESHOLD = 0.3

# Ollama Settings (Local AI)
OLLAMA_HOST = os.environ.get("OLLAMA_HOST", "http://localhost:11434")
PHI3_MODEL = "phi3:mini"  # Text processing
LLAVA_MODEL = "llava-phi3"  # Image processing

# Whisper Settings (Audio/Video)
WHISPER_MODEL = "tiny"  # Options: tiny, base, small, medium, large

# RAG Settings
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
MAX_RAG_RESULTS = 5

# Chatbot Settings
MAX_CONVERSATION_HISTORY = 3
ENABLE_RAG = True
'''

    # ============================================
    # FILE ORGANIZER - STATE
    # ============================================

    file_org_init = """# File Organizer package"""

    file_org_state = '''"""
Global State Management
"""

from threading import Lock


class FileOrganizerState:
    def __init__(self):
        self.watched_folders = []
        self.categories = []
        self.observers = {}
        self.lock = Lock()
        self.processing_queue = []
        self.ollama_client = None
        self.whisper_model = None
        self.is_initialized = False
        self.processing_stats = {
            "total": 0,
            "success": 0,
            "failed": 0,
            "by_type": {}
        }


state = FileOrganizerState()
'''

    # ============================================
    # FILE ORGANIZER - MODELS
    # ============================================

    file_org_models = '''"""
AI Model Initialization - Ollama + Whisper
"""

import ollama
from faster_whisper import WhisperModel
from config.settings import PHI3_MODEL, LLAVA_MODEL, WHISPER_MODEL, OLLAMA_HOST
from .state import state


def initialize_models():
    """Initialize Ollama and Whisper models"""
    print("🔄 Initializing AI models...")
    
    try:
        # Initialize Ollama client
        state.ollama_client = ollama.Client(host=OLLAMA_HOST)
        
        # Check if models are available
        print(f"📦 Checking Ollama models...")
        available_models = state.ollama_client.list()
        model_names = [m['name'] for m in available_models['models']]
        
        if PHI3_MODEL not in model_names:
            print(f"⚠️  {PHI3_MODEL} not found. Run: ollama pull {PHI3_MODEL}")
        else:
            print(f"✅ {PHI3_MODEL} ready")
        
        if LLAVA_MODEL not in model_names:
            print(f"⚠️  {LLAVA_MODEL} not found. Run: ollama pull {LLAVA_MODEL}")
        else:
            print(f"✅ {LLAVA_MODEL} ready")
        
        # Initialize Whisper
        print(f"🔄 Loading Whisper {WHISPER_MODEL} model...")
        state.whisper_model = WhisperModel(WHISPER_MODEL, device="cpu", compute_type="int8")
        print(f"✅ Whisper {WHISPER_MODEL} ready")
        
        state.is_initialized = True
        print("✅ All AI models initialized")
        
    except Exception as e:
        print(f"❌ Error initializing models: {e}")
        print("Make sure Ollama is running: ollama serve")
        state.is_initialized = False
'''

    # ============================================
    # FILE ORGANIZER - EXTRACTORS
    # ============================================

    file_org_extractors = '''"""
Multi-Modal File Extractors
Supports: Documents, Images, Audio, Video, Archives, Adobe files
"""

from pathlib import Path
from typing import Dict, List, Any
import json
import mimetypes

# Document extraction
import PyPDF2
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

# Image extraction
from PIL import Image
import pytesseract

# Audio/Video extraction
from faster_whisper import WhisperModel
import subprocess

# Adobe files
try:
    from psd_tools import PSDImage
    HAS_PSD = True
except ImportError:
    HAS_PSD = False

# Archives
import zipfile

from .state import state


class MultiModalExtractor:
    """Extract content from any file type"""
    
    @staticmethod
    def extract(file_path: str) -> Dict[str, Any]:
        """
        Extract all content from a file.
        Returns: {
            "text": str,
            "images": List[str],  # Base64 or paths
            "audio_transcript": str,
            "metadata": dict,
            "file_type": str
        }
        """
        path = Path(file_path)
        ext = path.suffix.lower()
        
        result = {
            "text": "",
            "images": [],
            "audio_transcript": "",
            "metadata": {
                "filename": path.name,
                "size": path.stat().st_size,
                "extension": ext
            },
            "file_type": MultiModalExtractor._detect_type(ext)
        }
        
        try:
            # Route to appropriate extractor
            if ext in ['.txt', '.md', '.json', '.xml', '.html', '.css', '.js', '.py', '.java', '.cpp', '.csv']:
                result["text"] = MultiModalExtractor._extract_text(file_path)
            
            elif ext == '.pdf':
                result.update(MultiModalExtractor._extract_pdf(file_path))
            
            elif ext in ['.docx', '.doc']:
                result.update(MultiModalExtractor._extract_docx(file_path))
            
            elif ext in ['.xlsx', '.xls']:
                result["text"] = MultiModalExtractor._extract_excel(file_path)
            
            elif ext in ['.pptx', '.ppt']:
                result.update(MultiModalExtractor._extract_pptx(file_path))
            
            elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                result.update(MultiModalExtractor._extract_image(file_path))
            
            elif ext == '.psd' and HAS_PSD:
                result.update(MultiModalExtractor._extract_psd(file_path))
            
            elif ext in ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.wma']:
                result["audio_transcript"] = MultiModalExtractor._extract_audio(file_path)
            
            elif ext in ['.mp4', '.avi', '.mov', '.mkv', '.webm']:
                result["audio_transcript"] = MultiModalExtractor._extract_video_audio(file_path)
            
            elif ext == '.zip':
                result["text"] = MultiModalExtractor._extract_zip(file_path)
            
            else:
                result["text"] = f"Filename: {path.stem}"
            
        except Exception as e:
            print(f"⚠️ Error extracting {file_path}: {e}")
            result["text"] = path.stem
            result["metadata"]["error"] = str(e)
        
        return result
    
    @staticmethod
    def _detect_type(ext: str) -> str:
        """Detect broad file category"""
        text_types = ['.txt', '.md', '.json', '.xml', '.html', '.css', '.js', '.py', '.java', '.csv']
        doc_types = ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt']
        image_types = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp', '.psd']
        audio_types = ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.wma']
        video_types = ['.mp4', '.avi', '.mov', '.mkv', '.webm']
        
        if ext in text_types or ext in doc_types:
            return "document"
        elif ext in image_types:
            return "image"
        elif ext in audio_types:
            return "audio"
        elif ext in video_types:
            return "video"
        else:
            return "other"
    
    @staticmethod
    def _extract_text(file_path: str) -> str:
        """Extract plain text"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    @staticmethod
    def _extract_pdf(file_path: str) -> Dict:
        """Extract text and images from PDF"""
        text = ""
        images = []
        
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
                text += "\\n"
        
        return {"text": text, "images": images}
    
    @staticmethod
    def _extract_docx(file_path: str) -> Dict:
        """Extract text and images from Word"""
        doc = DocxDocument(file_path)
        text = "\\n".join([p.text for p in doc.paragraphs])
        
        # TODO: Extract embedded images
        images = []
        
        return {"text": text, "images": images}
    
    @staticmethod
    def _extract_excel(file_path: str) -> str:
        """Extract data from Excel"""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"\\n=== Sheet: {sheet_name} ===\\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                text += row_text + "\\n"
        
        return text
    
    @staticmethod
    def _extract_pptx(file_path: str) -> Dict:
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        text = ""
        images = []
        
        for i, slide in enumerate(prs.slides):
            text += f"\\n=== Slide {i+1} ===\\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\\n"
        
        return {"text": text, "images": images}
    
    @staticmethod
    def _extract_image(file_path: str) -> Dict:
        """Extract text from image using OCR"""
        img = Image.open(file_path)
        
        # OCR for text
        try:
            text = pytesseract.image_to_string(img)
        except:
            text = ""
        
        # Store image path for vision model
        images = [file_path]
        
        return {"text": text, "images": images}
    
    @staticmethod
    def _extract_psd(file_path: str) -> Dict:
        """Extract preview from Photoshop"""
        psd = PSDImage.open(file_path)
        
        # Get layer names as text
        text = "Layers: " + ", ".join([layer.name for layer in psd])
        
        # TODO: Export composite image for vision processing
        images = []
        
        return {"text": text, "images": images}
    
    @staticmethod
    def _extract_audio(file_path: str) -> str:
        """Transcribe audio using Whisper"""
        if not state.whisper_model:
            return ""
        
        try:
            segments, info = state.whisper_model.transcribe(file_path, beam_size=5)
            transcript = " ".join([segment.text for segment in segments])
            return transcript
        except Exception as e:
            print(f"⚠️ Whisper error: {e}")
            return ""
    
    @staticmethod
    def _extract_video_audio(file_path: str) -> str:
        """Extract and transcribe audio from video"""
        # Extract audio using ffmpeg
        import tempfile
        
        temp_audio = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
        temp_audio_path = temp_audio.name
        temp_audio.close()
        
        try:
            # Extract audio track
            subprocess.run([
                "ffmpeg", "-i", file_path,
                "-vn", "-acodec", "pcm_s16le",
                "-ar", "16000", "-ac", "1",
                temp_audio_path
            ], capture_output=True, check=True)
            
            # Transcribe
            transcript = MultiModalExtractor._extract_audio(temp_audio_path)
            
            # Cleanup
            Path(temp_audio_path).unlink()
            
            return transcript
        except Exception as e:
            print(f"⚠️ Video audio extraction error: {e}")
            return ""
    
    @staticmethod
    def _extract_zip(file_path: str) -> str:
        """List contents of ZIP"""
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            file_list = zip_ref.namelist()
            return "Archive contents:\\n" + "\\n".join(file_list)
'''

    # ============================================
    # FILE ORGANIZER - AI PROCESSOR
    # ============================================

    file_org_ai_processor = '''"""
AI Processing with Ollama (Phi-3-Mini + LLaVA-Phi3)
"""

import ollama
from pathlib import Path
from typing import Dict, List
from config.settings import PHI3_MODEL, LLAVA_MODEL
from .state import state


class AIProcessor:
    """Process extracted content with AI models"""
    
    @staticmethod
    def analyze_content(extracted_data: Dict) -> Dict:
        """
        Analyze extracted content using appropriate AI models.
        Returns: {
            "summary": str,
            "keywords": List[str],
            "category_suggestions": List[str],
            "entities": List[str]
        }
        """
        result = {
            "summary": "",
            "keywords": [],
            "category_suggestions": [],
            "entities": []
        }
        
        # Process text with Phi-3-Mini
        if extracted_data.get("text"):
            text_analysis = AIProcessor._analyze_text(extracted_data["text"])
            result.update(text_analysis)
        
        # Process audio transcript
        if extracted_data.get("audio_transcript"):
            audio_analysis = AIProcessor._analyze_text(
                f"Audio transcript: {extracted_data['audio_transcript']}"
            )
            result["summary"] += "\\n\\nAudio: " + audio_analysis.get("summary", "")
        
        # Process images with LLaVA-Phi3
        if extracted_data.get("images"):
            for img_path in extracted_data["images"]:
                img_analysis = AIProcessor._analyze_image(img_path)
                result["summary"] += "\\n\\nImage: " + img_analysis
        
        return result
    
    @staticmethod
    def _analyze_text(text: str) -> Dict:
        """Analyze text with Phi-3-Mini"""
        if not state.ollama_client or not text.strip():
            return {"summary": "", "keywords": [], "category_suggestions": []}
        
        try:
            # Truncate very long text
            if len(text) > 4000:
                text = text[:4000] + "..."
            
            prompt = f"""Analyze this document and provide:
1. A brief summary (2-3 sentences)
2. 5 key keywords
3. 3 category suggestions for organizing this file

Document:
{text}

Respond in this format:
SUMMARY: <summary>
KEYWORDS: <keyword1>, <keyword2>, ...
CATEGORIES: <category1>, <category2>, <category3>
"""
            
            response = state.ollama_client.generate(
                model=PHI3_MODEL,
                prompt=prompt
            )
            
            output = response['response']
            
            # Parse response
            result = {
                "summary": "",
                "keywords": [],
                "category_suggestions": []
            }
            
            for line in output.split('\\n'):
                if line.startswith('SUMMARY:'):
                    result["summary"] = line.replace('SUMMARY:', '').strip()
                elif line.startswith('KEYWORDS:'):
                    keywords = line.replace('KEYWORDS:', '').strip()
                    result["keywords"] = [k.strip() for k in keywords.split(',')]
                elif line.startswith('CATEGORIES:'):
                    categories = line.replace('CATEGORIES:', '').strip()
                    result["category_suggestions"] = [c.strip() for c in categories.split(',')]
            
            return result
            
        except Exception as e:
            print(f"⚠️ Phi-3-Mini error: {e}")
            return {"summary": "", "keywords": [], "category_suggestions": []}
    
    @staticmethod
    def _analyze_image(image_path: str) -> str:
        """Analyze image with LLaVA-Phi3"""
        if not state.ollama_client:
            return ""
        
        try:
            response = state.ollama_client.generate(
                model=LLAVA_MODEL,
                prompt="Describe this image in detail. What is the main subject? What text or important elements are visible?",
                images=[image_path]
            )
            
            return response['response']
            
        except Exception as e:
            print(f"⚠️ LLaVA-Phi3 error: {e}")
            return ""
    
    @staticmethod
    def classify_into_categories(analysis: Dict, categories: List[Dict]) -> Dict:
        """
        Use AI to classify into user-defined categories.
        Returns: {"category": str, "confidence": float}
        """
        if not categories or not state.ollama_client:
            return {"category": None, "confidence": 0.0}
        
        try:
            category_names = [c['name'] for c in categories]
            
            prompt = f"""Given this file analysis, which category fits best?

Summary: {analysis.get('summary', '')}
Keywords: {', '.join(analysis.get('keywords', []))}

Available categories: {', '.join(category_names)}

Respond with ONLY the category name that fits best."""
            
            response = state.ollama_client.generate(
                model=PHI3_MODEL,
                prompt=prompt
            )
            
            predicted_category = response['response'].strip()
            
            # Find best match
            for cat in categories:
                if cat['name'].lower() in predicted_category.lower():
                    return {"category": cat['name'], "confidence": 0.85}
            
            # Fallback to first suggestion
            if analysis.get('category_suggestions'):
                for suggestion in analysis['category_suggestions']:
                    for cat in categories:
                        if cat['name'].lower() in suggestion.lower():
                            return {"category": cat['name'], "confidence": 0.70}
            
            return {"category": None, "confidence": 0.0}
            
        except Exception as e:
            print(f"⚠️ Classification error: {e}")
            return {"category": None, "confidence": 0.0}
'''

    # ============================================
    # FILE ORGANIZER - RAG GENERATOR
    # ============================================

    file_org_rag = '''"""
RAG (Retrieval Augmented Generation) System
Stores processed files in searchable format for chatbot
"""

import json
from pathlib import Path
from datetime import datetime
from typing import List, Dict
from config.settings import RAG_DIR, CHUNK_SIZE, CHUNK_OVERLAP


class RAGGenerator:
    """Generate and store RAG data for files"""
    
    @staticmethod
    def create_rag_document(
        file_path: str,
        extracted_data: Dict,
        ai_analysis: Dict,
        category: str
    ) -> str:
        """
        Create a RAG document with all file information.
        Returns: Path to RAG JSON file
        """
        file_name = Path(file_path).stem
        rag_filename = f"{file_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.rag.json"
        rag_path = RAG_DIR / rag_filename
        
        # Chunk the text content
        full_text = extracted_data.get("text", "")
        if extracted_data.get("audio_transcript"):
            full_text += "\\n\\nAudio Transcript:\\n" + extracted_data["audio_transcript"]
        
        chunks = RAGGenerator._chunk_text(full_text)
        
        # Create RAG document
        rag_doc = {
            "file_info": {
                "original_path": file_path,
                "filename": Path(file_path).name,
                "category": category,
                "file_type": extracted_data.get("file_type", "unknown"),
                "created_at": datetime.now().isoformat()
            },
            "content": {
                "full_text": full_text[:10000],  # Store first 10k chars
                "chunks": chunks,
                "audio_transcript": extracted_data.get("audio_transcript", ""),
                "images": extracted_data.get("images", [])
            },
            "analysis": {
                "summary": ai_analysis.get("summary", ""),
                "keywords": ai_analysis.get("keywords", []),
                "entities": ai_analysis.get("entities", []),
                "category_suggestions": ai_analysis.get("category_suggestions", [])
            },
            "metadata": extracted_data.get("metadata", {})
        }
        
        # Save RAG document
        with open(rag_path, 'w', encoding='utf-8') as f:
            json.dump(rag_doc, f, indent=2, ensure_ascii=False)
        
        print(f"💾 RAG document created: {rag_path}")
        return str(rag_path)
    
    @staticmethod
    def _chunk_text(text: str) -> List[str]:
        """Split text into overlapping chunks"""
        if not text:
            return []
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + CHUNK_SIZE
            chunk = text[start:end]
            chunks.append(chunk)
            start += CHUNK_SIZE - CHUNK_OVERLAP
        
        return chunks
    
    @staticmethod
    def search_rag(query: str, max_results: int = 5) -> List[Dict]:
        """
        Search through all RAG documents for relevant information.
        Returns list of relevant documents with scores.
        """
        results = []
        query_lower = query.lower()
        query_words = set(query_lower.split())
        
        # Search through all RAG files
        for rag_file in RAG_DIR.glob("*.rag.json"):
            try:
                with open(rag_file, 'r', encoding='utf-8') as f:
                    rag_doc = json.load(f)
                
                # Calculate relevance score
                score = 0
                
                # Check summary
                summary = rag_doc.get("analysis", {}).get("summary", "").lower()
                score += sum(word in summary for word in query_words) * 3
                
                # Check keywords
                keywords = [k.lower() for k in rag_doc.get("analysis", {}).get("keywords", [])]
                score += sum(word in keywords for word in query_words) * 5
                
                # Check chunks
                for chunk in rag_doc.get("content", {}).get("chunks", []):
                    chunk_lower = chunk.lower()
                    score += sum(word in chunk_lower for word in query_words)
                
                # Check audio transcript
                transcript = rag_doc.get("content", {}).get("audio_transcript", "").lower()
                score += sum(word in transcript for word in query_words) * 2
                
                if score > 0:
                    results.append({
                        "file": rag_doc["file_info"]["filename"],
                        "path": rag_doc["file_info"]["original_path"],
                        "summary": rag_doc["analysis"]["summary"],
                        "keywords": rag_doc["analysis"]["keywords"],
                        "score": score,
                        "content_preview": rag_doc["content"]["full_text"][:500]
                    })
            
            except Exception as e:
                print(f"⚠️ Error reading RAG file {rag_file}: {e}")
                continue
        
        # Sort by score and return top results
        results.sort(key=lambda x: x["score"], reverse=True)
        return results[:max_results]
'''

    # ============================================
    # FILE ORGANIZER - FILE PROCESSOR
    # ============================================

    file_org_processor = '''"""
File Processing Pipeline with RAG Generation
"""

import shutil
from pathlib import Path
from config.settings import CONFIDENCE_THRESHOLD
from .state import state
from .file_extractors import MultiModalExtractor
from .ai_processor import AIProcessor
from .rag_generator import RAGGenerator


def process_file(file_path: str, folder_id: int):
    """
    Complete file processing pipeline:
    1. Extract content (text, images, audio)
    2. Analyze with AI
    3. Classify into category
    4. Generate RAG document
    5. Move to destination
    """
    try:
        print(f"\\n📄 Processing: {file_path}")
        
        # Step 1: Extract all content
        print("  🔍 Extracting content...")
        extracted_data = MultiModalExtractor.extract(file_path)
        
        # Step 2: Analyze with AI
        print("  🤖 Analyzing with AI...")
        ai_analysis = AIProcessor.analyze_content(extracted_data)
        
        # Step 3: Classify into category
        print("  🏷️  Classifying...")
        classification = AIProcessor.classify_into_categories(ai_analysis, state.categories)
        category_name = classification["category"]
        confidence = classification["confidence"]
        
        # Check if we have a valid category
        if category_name and confidence >= CONFIDENCE_THRESHOLD:
            category = next((c for c in state.categories if c['name'] == category_name), None)
            
            if category:
                # Step 4: Generate RAG document BEFORE moving
                print("  💾 Generating RAG document...")
                rag_path = RAGGenerator.create_rag_document(
                    file_path,
                    extracted_data,
                    ai_analysis,
                    category_name
                )
                
                # Step 5: Move file to destination
                dest_folder = Path(category['path'])
                dest_folder.mkdir(parents=True, exist_ok=True)
                dest_path = dest_folder / Path(file_path).name
                
                # Handle duplicates
                counter = 1
                while dest_path.exists():
                    dest_path = dest_folder / f"{Path(file_path).stem}_{counter}{Path(file_path).suffix}"
                    counter += 1
                
                shutil.move(file_path, dest_path)
                print(f"  ✅ Moved to: {dest_path}")
                print(f"  📊 Confidence: {confidence:.2%}")
                
                # Update stats
                state.processing_stats["total"] += 1
                state.processing_stats["success"] += 1
                file_type = extracted_data.get("file_type", "other")
                state.processing_stats["by_type"][file_type] = state.processing_stats["by_type"].get(file_type, 0) + 1
                
                return {
                    "status": "success",
                    "file_name": dest_path.name,
                    "category": category_name,
                    "confidence": confidence,
                    "summary": ai_analysis.get("summary", ""),
                    "keywords": ai_analysis.get("keywords", []),
                    "rag_path": rag_path
                }
        
        # Low confidence - generate RAG anyway but don't move
        print(f"  ⚠️ Low confidence ({confidence:.2%})")
        rag_path = RAGGenerator.create_rag_document(
            file_path,
            extracted_data,
            ai_analysis,
            "Uncategorized"
        )
        
        state.processing_stats["total"] += 1
        
        return {
            "status": "low_confidence",
            "file_name": Path(file_path).name,
            "confidence": confidence,
            "summary": ai_analysis.get("summary", ""),
            "suggestions": ai_analysis.get("category_suggestions", []),
            "rag_path": rag_path
        }
    
    except Exception as e:
        print(f"  ❌ Error: {e}")
        import traceback
        traceback.print_exc()
        
        state.processing_stats["total"] += 1
        state.processing_stats["failed"] += 1
        
        return {
            "status": "error",
            "file_name": Path(file_path).name,
            "error": str(e)
        }
'''

    # ============================================
    # FILE ORGANIZER - FILE WATCHER (MISSING FILE!)
    # ============================================

    file_org_watcher = '''"""
File System Watcher
Monitors folders for new files and processes them automatically
"""

import time
from pathlib import Path
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from .state import state
from .file_processor import process_file


class FileHandler(FileSystemEventHandler):
    """Handle file system events"""
    
    def __init__(self, folder_id: int):
        self.folder_id = folder_id
        self.processing_delay = 1.0  # Wait 1 second before processing
        self.pending_files = {}
    
    def on_created(self, event):
        """Handle file creation"""
        if event.is_directory:
            return
        
        file_path = event.src_path
        print(f"🆕 New file detected: {file_path}")
        
        # Add to pending queue with timestamp
        self.pending_files[file_path] = time.time()
    
    def on_modified(self, event):
        """Handle file modification"""
        if event.is_directory:
            return
        
        file_path = event.src_path
        # Update timestamp if already pending
        if file_path in self.pending_files:
            self.pending_files[file_path] = time.time()
    
    def process_pending_files(self):
        """Process files that haven't been modified for delay period"""
        current_time = time.time()
        files_to_process = []
        
        for file_path, timestamp in list(self.pending_files.items()):
            if current_time - timestamp >= self.processing_delay:
                files_to_process.append(file_path)
                del self.pending_files[file_path]
        
        for file_path in files_to_process:
            if Path(file_path).exists():
                try:
                    process_file(file_path, self.folder_id)
                except Exception as e:
                    print(f"❌ Error processing {file_path}: {e}")


def start_watching(folder_id: int, folder_path: str):
    """Start watching a folder for new files"""
    if folder_id in state.observers:
        print(f"⚠️ Already watching folder {folder_id}")
        return
    
    if not Path(folder_path).exists():
        print(f"❌ Folder does not exist: {folder_path}")
        return
    
    event_handler = FileHandler(folder_id)
    observer = Observer()
    observer.schedule(event_handler, folder_path, recursive=False)
    observer.start()
    
    state.observers[folder_id] = {
        "observer": observer,
        "handler": event_handler
    }
    
    print(f"👁️  Watching folder: {folder_path}")
    
    # Start processing loop in background
    def processing_loop():
        while folder_id in state.observers:
            event_handler.process_pending_files()
            time.sleep(0.5)
    
    from threading import Thread
    thread = Thread(target=processing_loop, daemon=True)
    thread.start()


def stop_watching(folder_id: int):
    """Stop watching a folder"""
    if folder_id not in state.observers:
        return
    
    observer_data = state.observers[folder_id]
    observer_data["observer"].stop()
    observer_data["observer"].join()
    
    del state.observers[folder_id]
    print(f"👁️‍🗨️ Stopped watching folder {folder_id}")


def start_all_watchers():
    """Start watching all active folders"""
    for folder in state.watched_folders:
        if folder.get("status") == "Active":
            start_watching(folder["id"], folder["path"])
'''

    # ============================================
    # FILE ORGANIZER - ROUTES
    # ============================================

    file_org_routes = '''"""
API Routes for File Organizer
"""

import os
import json
import time
from pathlib import Path
from flask import request, jsonify
from config.settings import CONFIG_FILE
from .state import state
from .file_watcher import start_watching, stop_watching
from .file_processor import process_file
from .rag_generator import RAGGenerator


def load_config():
    """Load watched folders and categories from disk"""
    if os.path.exists(CONFIG_FILE):
        with open(CONFIG_FILE, "r") as f:
            config = json.load(f)
            state.watched_folders = config.get("watched_folders", [])
            state.categories = config.get("categories", [])
            print(f"📂 Loaded {len(state.watched_folders)} watched folders")
            print(f"🏷️  Loaded {len(state.categories)} categories")


def save_config():
    """Save watched folders and categories"""
    config = {
        "watched_folders": state.watched_folders,
        "categories": state.categories
    }
    with open(CONFIG_FILE, "w") as f:
        json.dump(config, f, indent=2)


def register_routes(app):
    """Register all API routes"""
    
    @app.route("/api/health", methods=["GET"])
    def health():
        return jsonify({
            "status": "online",
            "initialized": state.is_initialized,
            "stats": state.processing_stats
        })

    # Watched Folders
    @app.route("/api/watched-folders", methods=["GET"])
    def get_folders():
        return jsonify(state.watched_folders)

    @app.route("/api/watched-folders", methods=["POST"])
    def add_folder():
        data = request.json
        folder = {
            "id": int(time.time() * 1000),
            "name": data.get("name"),
            "path": data.get("path"),
            "status": "Active"
        }
        state.watched_folders.append(folder)
        save_config()
        start_watching(folder["id"], folder["path"])
        return jsonify(folder)

    @app.route("/api/watched-folders/<int:folder_id>/toggle", methods=["POST"])
    def toggle_folder(folder_id):
        folder = next((f for f in state.watched_folders if f["id"] == folder_id), None)
        if folder:
            folder["status"] = "Paused" if folder["status"] == "Active" else "Active"
            if folder["status"] == "Active":
                start_watching(folder["id"], folder["path"])
            else:
                stop_watching(folder["id"])
            save_config()
            return jsonify(folder)
        return jsonify({"error": "Folder not found"}), 404

    # Categories
    @app.route("/api/categories", methods=["GET"])
    def get_categories():
        return jsonify(state.categories)

    @app.route("/api/categories", methods=["POST"])
    def add_category():
        data = request.json
        category = {
            "id": int(time.time() * 1000),
            "name": data.get("name"),
            "path": data.get("path")
        }
        state.categories.append(category)
        save_config()
        return jsonify(category)

    # Manual folder processing
    @app.route("/api/process-folder/<int:folder_id>", methods=["POST"])
    def process_folder_route(folder_id):
        folder = next((f for f in state.watched_folders if f["id"] == folder_id), None)
        if not folder or not Path(folder["path"]).exists():
            return jsonify({"error": "Folder not found"}), 404
        results = []
        for file_path in Path(folder["path"]).glob("*"):
            if file_path.is_file():
                results.append(process_file(str(file_path), folder_id))
        return jsonify({"processed": len(results), "results": results})
    
    # RAG Search (for chatbot integration)
    @app.route("/api/rag/search", methods=["POST"])
    def search_rag():
        query = request.json.get("query", "")
        max_results = request.json.get("max_results", 5)
        
        if not query:
            return jsonify({"error": "Query required"}), 400
        
        results = RAGGenerator.search_rag(query, max_results)
        return jsonify({"results": results})
    
    # Get processing statistics
    @app.route("/api/stats", methods=["GET"])
    def get_stats():
        return jsonify(state.processing_stats)
'''

    # ============================================
    # FILE ORGANIZER - APP (MISSING FILE!)
    # ============================================

    file_org_app = '''"""
File Organizer Flask Application
"""

from flask import Flask
from flask_cors import CORS
from .models import initialize_models
from .routes import register_routes, load_config
from .file_watcher import start_all_watchers


def create_file_organizer_app():
    """Create and configure the file organizer Flask app"""
    app = Flask(__name__)
    CORS(app)
    
    # Initialize AI models
    initialize_models()
    
    # Load configuration
    load_config()
    
    # Register API routes
    register_routes(app)
    
    # Start file watchers
    start_all_watchers()
    
    print("📡 File Organizer API ready")
    print("Endpoints:")
    print("  GET  /api/health")
    print("  GET  /api/watched-folders")
    print("  POST /api/watched-folders")
    print("  POST /api/watched-folders/<id>/toggle")
    print("  GET  /api/categories")
    print("  POST /api/categories")
    print("  POST /api/process-folder/<id>")
    print("  POST /api/rag/search")
    print("  GET  /api/stats")
    
    return app
'''

    # ============================================
    # CHATBOT - INIT
    # ============================================

    chatbot_init = """# Chatbot package"""

    # ============================================
    # CHATBOT - OLLAMA CLIENT
    # ============================================

    chatbot_ollama = '''"""
Ollama Chatbot with RAG Integration
"""

import ollama
import requests
from config.settings import (
    OLLAMA_HOST, PHI3_MODEL, 
    MAX_CONVERSATION_HISTORY, 
    ENABLE_RAG, FILE_ORGANIZER_PORT
)


class OllamaChatbot:
    def __init__(self):
        self.client = ollama.Client(host=OLLAMA_HOST)
        self.conversation_history = []
        self.is_ready = True
        self.loading = False
        self.system_instruction = """You are a helpful AI assistant for Leologic, an AI-powered file organization system.

You help users:
1. Organize and find their files
2. Answer questions about file contents
3. Provide summaries of documents
4. Search through organized files

When users ask about their files, you can search through them using the RAG system."""
    
    def generate_response(self, user_message, store_history=True):
        """Generate response with optional RAG context"""
        try:
            messages = [{"role": "system", "content": self.system_instruction}]
            
            # Add conversation history
            for msg in self.conversation_history[-MAX_CONVERSATION_HISTORY:]:
                messages.append({"role": "user", "content": msg["user"]})
                messages.append({"role": "assistant", "content": msg["assistant"]})
            
            # Check if we should use RAG
            rag_context = ""
            if ENABLE_RAG and self._should_use_rag(user_message):
                rag_context = self._get_rag_context(user_message)
                if rag_context:
                    messages.append({
                        "role": "system",
                        "content": f"Relevant file information:\\n{rag_context}"
                    })
            
            # Add current user message
            messages.append({"role": "user", "content": user_message})
            
            # Generate response
            response = self.client.chat(
                model=PHI3_MODEL,
                messages=messages
            )
            
            response_text = response['message']['content']
            
            # Save to history
            if store_history:
                self.conversation_history.append({
                    "user": user_message,
                    "assistant": response_text
                })
            
            return response_text
        
        except Exception as e:
            print(f"❌ Error generating response: {e}")
            import traceback
            traceback.print_exc()
            return f"Error: {str(e)}"
    
    def _should_use_rag(self, message: str) -> bool:
        """Determine if query needs file search"""
        rag_keywords = [
            'file', 'document', 'find', 'search', 'what did', 
            'show me', 'where is', 'tell me about', 'summary',
            'contents', 'pdf', 'word', 'excel', 'image', 'video'
        ]
        message_lower = message.lower()
        return any(keyword in message_lower for keyword in rag_keywords)
    
    def _get_rag_context(self, query: str) -> str:
        """Get relevant context from RAG system"""
        try:
            # Call file organizer's RAG search endpoint
            response = requests.post(
                f"http://localhost:{FILE_ORGANIZER_PORT}/api/rag/search",
                json={"query": query, "max_results": 3}
            )
            
            if response.status_code == 200:
                results = response.json().get("results", [])
                
                if not results:
                    return ""
                
                # Format results as context
                context = "Here are relevant files:\\n\\n"
                for i, result in enumerate(results, 1):
                    context += f"{i}. **{result['file']}**\\n"
                    context += f"   Summary: {result['summary']}\\n"
                    context += f"   Keywords: {', '.join(result['keywords'])}\\n"
                    context += f"   Preview: {result['content_preview'][:200]}...\\n\\n"
                
                return context
            
        except Exception as e:
            print(f"⚠️ RAG search error: {e}")
        
        return ""
    
    def clear_history(self):
        """Clear conversation history"""
        self.conversation_history = []
'''

    # ============================================
    # CHATBOT - ROUTES
    # ============================================

    chatbot_routes = '''"""
API Routes for Chatbot
"""

import time
from flask import request, jsonify


def register_routes(app, chatbot):
    """Register all chatbot API routes"""
    
    @app.route("/health", methods=["GET"])
    def health_check():
        return jsonify({
            "status": "online",
            "model_ready": chatbot.is_ready,
            "loading": chatbot.loading
        })

    @app.route("/chat", methods=["POST"])
    def chat():
        data = request.json
        user_message = data.get("message", "").strip()

        if not user_message:
            return jsonify({"error": "Message is required"}), 400

        if not chatbot.is_ready:
            return jsonify({
                "response": "Ollama is still initializing...",
                "model_ready": False
            })

        start_time = time.time()
        response = chatbot.generate_response(user_message)
        elapsed = round(time.time() - start_time, 2)

        return jsonify({
            "response": response,
            "model_ready": chatbot.is_ready,
            "elapsed_time": elapsed
        })

    @app.route("/clear", methods=["POST"])
    def clear_history():
        chatbot.clear_history()
        return jsonify({
            "status": "success",
            "message": "Conversation history cleared"
        })

    @app.route("/status", methods=["GET"])
    def get_status():
        return jsonify({
            "conversation_length": len(chatbot.conversation_history),
            "model_ready": chatbot.is_ready
        })
'''

    # ============================================
    # CHATBOT - APP
    # ============================================

    chatbot_app = '''"""
Chatbot Flask Application
"""

from flask import Flask
from flask_cors import CORS
from .ollama_client import OllamaChatbot
from .routes import register_routes


def create_chatbot_app():
    """Create and configure the chatbot Flask app"""
    app = Flask(__name__)
    CORS(app)
    
    # Initialize chatbot
    chatbot = OllamaChatbot()
    
    # Register API routes
    register_routes(app, chatbot)
    
    print("📡 Chatbot API ready with RAG integration")
    print("Endpoints:")
    print("  GET  /health")
    print("  GET  /status")
    print("  POST /chat")
    print("  POST /clear")
    
    return app
'''

    # ============================================
    # UTILS - INIT (MISSING FILE!)
    # ============================================

    utils_init = """# Utils package"""

    # ============================================
    # UTILS - LOGGER (MISSING FILE!)
    # ============================================

    utils_logger = '''"""
Logging Utilities
"""

import logging
from pathlib import Path
from datetime import datetime


def setup_logger(name: str, log_file: str = None) -> logging.Logger:
    """
    Setup a logger with file and console handlers
    
    Args:
        name: Logger name
        log_file: Optional log file path
    
    Returns:
        Configured logger instance
    """
    logger = logging.getLogger(name)
    logger.setLevel(logging.INFO)
    
    # Avoid duplicate handlers
    if logger.handlers:
        return logger
    
    # Create formatter
    formatter = logging.Formatter(
        '%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    
    # Console handler
    console_handler = logging.StreamHandler()
    console_handler.setLevel(logging.INFO)
    console_handler.setFormatter(formatter)
    logger.addHandler(console_handler)
    
    # File handler (if specified)
    if log_file:
        log_path = Path(log_file)
        log_path.parent.mkdir(parents=True, exist_ok=True)
        
        file_handler = logging.FileHandler(log_file)
        file_handler.setLevel(logging.DEBUG)
        file_handler.setFormatter(formatter)
        logger.addHandler(file_handler)
    
    return logger


def log_processing(file_path: str, status: str, details: dict = None):
    """
    Log file processing events
    
    Args:
        file_path: Path to processed file
        status: Processing status (success, error, skipped)
        details: Additional details dictionary
    """
    logger = logging.getLogger("file_processor")
    
    message = f"File: {file_path} | Status: {status}"
    if details:
        message += f" | Details: {details}"
    
    if status == "error":
        logger.error(message)
    elif status == "success":
        logger.info(message)
    else:
        logger.warning(message)
'''

    # ============================================
    # REQUIREMENTS.TXT
    # ============================================

    requirements = """# Leologic Backend Requirements - Local AI Edition

# Flask & Web
flask==3.0.0
flask-cors==4.0.0
python-dotenv==1.0.0

# File Monitoring
watchdog==3.0.0

# Document Processing
PyPDF2==3.0.1
python-docx==1.1.0
openpyxl==3.1.2
python-pptx==0.6.23
pytesseract==0.3.10
Pillow==10.1.0

# Image Processing
psd-tools==1.9.27

# Audio/Video Processing
faster-whisper==1.0.0

# Local AI
ollama==0.1.6

# Utilities
requests==2.31.0
"""

    # ============================================
    # ENV EXAMPLE
    # ============================================

    env_example = """# Leologic Backend Environment Variables

# Ollama Configuration
OLLAMA_HOST=http://localhost:11434

# Optional: Custom paths
# DATA_DIR=/path/to/data
# RAG_DIR=/path/to/rag_storage
"""

    # ============================================
    # GITIGNORE
    # ============================================

    gitignore = """# Python
__pycache__/
*.py[cod]
*$py.class
*.so
.Python
build/
develop-eggs/
dist/
downloads/
eggs/
.eggs/
lib/
lib64/
parts/
sdist/
var/
wheels/
*.egg-info/
.installed.cfg
*.egg

# Virtual Environment
venv/
env/
ENV/

# IDE
.vscode/
.idea/
*.swp
*.swo
*~

# Data & RAG
data/
*.json
!requirements.txt
!package.json

# Environment
.env

# Logs
*.log
logs/
"""

    # ============================================
    # README
    # ============================================

    readme = """# Leologic Backend - Local AI Edition

AI-powered file organizer with local AI models (Ollama + Whisper) and RAG-enabled chatbot.

## 🚀 Quick Start

### 1. Install System Dependencies

**Ollama**: Visit https://ollama.ai  
**FFmpeg**: For video processing  
**Tesseract**: For OCR

### 2. Pull AI Models

```bash
ollama pull phi3:mini
ollama pull llava-phi3
```

### 3. Setup Python

```bash
cd backend
python -m venv venv
source venv/bin/activate  # or venv\\Scripts\\activate on Windows
pip install -r requirements.txt
```

### 4. Run

```bash
python main.py
```

Services start on:
- File Organizer: http://localhost:5001
- Chatbot: http://localhost:5000

See full documentation in the code!
"""

    # ============================================
    # CREATE ALL FILES
    # ============================================

    files_to_create = {
        # Root files
        "main.py": main_py,
        "requirements.txt": requirements,
        ".env.example": env_example,
        ".gitignore": gitignore,
        "README.md": readme,
        # Config
        "config/__init__.py": config_init,
        "config/settings.py": config_settings,
        # File Organizer
        "file_organizer/__init__.py": file_org_init,
        "file_organizer/state.py": file_org_state,
        "file_organizer/models.py": file_org_models,
        "file_organizer/file_extractors.py": file_org_extractors,
        "file_organizer/ai_processor.py": file_org_ai_processor,
        "file_organizer/rag_generator.py": file_org_rag,
        "file_organizer/file_processor.py": file_org_processor,
        "file_organizer/file_watcher.py": file_org_watcher,
        "file_organizer/routes.py": file_org_routes,
        "file_organizer/app.py": file_org_app,
        # Chatbot
        "chatbot/__init__.py": chatbot_init,
        "chatbot/ollama_client.py": chatbot_ollama,
        "chatbot/routes.py": chatbot_routes,
        "chatbot/app.py": chatbot_app,
        # Utils
        "utils/__init__.py": utils_init,
        "utils/logger.py": utils_logger,
    }

    for file_path, content in files_to_create.items():
        full_path = backend_dir / file_path
        create_file(full_path, content)

    print()
    print("=" * 60)
    print("✨ Leologic Backend setup complete!")
    print("=" * 60)
    print()
    print("📋 Next steps:")
    print("1. Install Ollama: https://ollama.ai")
    print("2. Pull models: ollama pull phi3:mini && ollama pull llava-phi3")
    print("3. Install deps: pip install -r requirements.txt")
    print("4. Run: python main.py")
    print()


if __name__ == "__main__":
    setup_backend()
